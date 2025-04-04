import os
from PIL import Image
from pptx import Presentation
import pytesseract as pytesseract
import pyttsx3
import docx
import PyPDF2
from django.shortcuts import render
from django.http import JsonResponse
from django.conf import settings
from django.views.decorators.csrf import csrf_exempt
from django.core.files.storage import default_storage


# Media folder for audio storage
MEDIA_FOLDER = os.path.join(settings.BASE_DIR, "media")
os.makedirs(MEDIA_FOLDER, exist_ok=True)

# Initialize TTS Engine
engine = pyttsx3.init()

# Fetch available voices
voices = engine.getProperty("voices")
VOICE_MAP = {
    "default": None,  # Uses system default voice
    "male": 0,        # First available voice
    "female": 1,      # Second available voice (if available)
}

def home(request):
    """Render homepage with all generated files."""
    functionality = request.GET.get("functionality", "tts")
    files = [f for f in os.listdir(MEDIA_FOLDER) if f.endswith(".mp3")]
    return render(request, "index.html", {"functionality": functionality, "files": files})


@csrf_exempt
def text_to_speech(request):
    """Convert text input to speech and return updated file list."""
    if request.method == "POST":
        text = request.POST.get("text", "").strip()
        speed = float(request.POST.get("speed", 150))
        pitch = float(request.POST.get("pitch", 1.0))
        voice = request.POST.get("voice", "default")

        if not text:
            return JsonResponse({"error": "Text field is required."}, status=400)

        # Set voice properties
        engine.setProperty("rate", speed)
        engine.setProperty("pitch", pitch)
        # print(VOICE_MAP)
        # Assign voice if available
        if voice in VOICE_MAP and VOICE_MAP[voice] is not None and VOICE_MAP[voice] < len(voices):
            engine.setProperty("voice", voices[VOICE_MAP[voice]].id)

        file_name = f"speech_{len(os.listdir(MEDIA_FOLDER)) + 1}.mp3"
        file_path = os.path.join(MEDIA_FOLDER, file_name)

        engine.save_to_file(text, file_path)
        engine.runAndWait()

        files = [f for f in os.listdir(MEDIA_FOLDER) if f.endswith(".mp3")]
        return JsonResponse({"audio_url": f"/media/{file_name}", "files": files})


@csrf_exempt
def file_to_speech(request):
    """Convert uploaded PDF/DOCX to speech and return updated file list."""
    if request.method == "POST" and request.FILES.get("file"):
        uploaded_file = request.FILES["file"]
        file_extension = uploaded_file.name.split(".")[-1].lower()
        speed = float(request.POST.get("speed", 150))
        pitch = float(request.POST.get("pitch", 1.0))
        voice = request.POST.get("voice", "default")

        file_path = os.path.join(MEDIA_FOLDER, uploaded_file.name)
        with default_storage.open(file_path, "wb+") as destination:
            for chunk in uploaded_file.chunks():
                destination.write(chunk)

        extracted_text = ""

        try:
            if file_extension == "pdf":
                with open(file_path, "rb") as f:
                    pdf_reader = PyPDF2.PdfReader(f)
                    for page in pdf_reader.pages:
                        extracted_text += page.extract_text() + " "
            elif file_extension == "docx":
                doc = docx.Document(file_path)
                extracted_text = " ".join([para.text for para in doc.paragraphs])
            else:
                return JsonResponse({"error": "Unsupported file type."}, status=400)

            if not extracted_text.strip():
                return JsonResponse({"error": "No text found in file."}, status=400)

            # Set TTS properties
            engine.setProperty("rate", speed)
            engine.setProperty("pitch", pitch)
            print(VOICE_MAP)
            if voice in VOICE_MAP and VOICE_MAP[voice] is not None and VOICE_MAP[voice] < len(voices):
                engine.setProperty("voice", voices[VOICE_MAP[voice]].id)

            audio_filename = f"file_speech_{len(os.listdir(MEDIA_FOLDER)) + 1}.mp3"
            audio_filepath = os.path.join(MEDIA_FOLDER, audio_filename)

            engine.save_to_file(extracted_text, audio_filepath)
            engine.runAndWait()

            files = [f for f in os.listdir(MEDIA_FOLDER) if f.endswith(".mp3")]
            return JsonResponse({"audio_url": f"/media/{audio_filename}", "files": files})

        except Exception as e:
            return JsonResponse({"error": f"Error: {str(e)}"}, status=500)

    return JsonResponse({"error": "Invalid request."}, status=400)


@csrf_exempt
def image_to_speech(request):
    """Extract text from an image and convert it to speech."""
    if request.method == "POST" and request.FILES.get("image"):
        image_file = request.FILES["image"]
        file_extension = image_file.name.split(".")[-1].lower()

        # Save uploaded image
        image_path = os.path.join(MEDIA_FOLDER, image_file.name)
        with default_storage.open(image_path, "wb+") as destination:
            for chunk in image_file.chunks():
                destination.write(chunk)

        try:
            # Open and process the image
            image = Image.open(image_path)
            extracted_text = pytesseract.image_to_string(image)

            if not extracted_text.strip():
                return JsonResponse({"error": "No readable text found in the image."}, status=400)

            # Generate speech file
            audio_filename = f"image_speech_{len(os.listdir(MEDIA_FOLDER)) + 1}.mp3"
            audio_filepath = os.path.join(MEDIA_FOLDER, audio_filename)

            engine.save_to_file(extracted_text, audio_filepath)
            engine.runAndWait()
            print(audio_filepath)
            return JsonResponse({"audio_url": f"/media/{audio_filename}"})

        except Exception as e:
            return JsonResponse({"error": f"Error: {str(e)}"}, status=500)

    return JsonResponse({"error": "Invalid request."}, status=400)


@csrf_exempt
def pptx_to_speech(request):
    """Extract text from a PPTX file and convert it to speech."""
    if request.method == "POST" and request.FILES.get("ppt"):
        pptx_file = request.FILES["ppt"]
        speed = float(request.POST.get("speed", 150))
        pitch = float(request.POST.get("pitch", 1.0))
        voice = request.POST.get("voice", "default")

        if not pptx_file.name.lower().endswith(".pptx"):
            return JsonResponse({"error": "Only .pptx files are supported."}, status=400)

        # Save the uploaded PPTX file
        file_path = os.path.join(MEDIA_FOLDER, pptx_file.name)
        with default_storage.open(file_path, "wb+") as destination:
            for chunk in pptx_file.chunks():
                destination.write(chunk)
        print("File written")
        try:
            # Parse presentation and extract text
            prs = Presentation(file_path)
            text_runs = []
            print("Entering for loop")
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame") and shape.text_frame is not None:
                        for paragraph in shape.text_frame.paragraphs:
                            full_line = ''.join([run.text for run in paragraph.runs]).strip()
                            if full_line:
                                text_runs.append(full_line)

            extracted_text = "\n".join(text_runs)
            print("Extracted text")

            if not extracted_text.strip():
                return JsonResponse({"error": "No readable text found in the presentation."}, status=400)

            # Set TTS engine properties
            engine.setProperty("rate", speed)
            engine.setProperty("pitch", pitch)
            if voice in VOICE_MAP and VOICE_MAP[voice] is not None and VOICE_MAP[voice] < len(voices):
                engine.setProperty("voice", voices[VOICE_MAP[voice]].id)

            audio_filename = f"pptx_speech_{len(os.listdir(MEDIA_FOLDER)) + 1}.mp3"
            audio_filepath = os.path.join(MEDIA_FOLDER, audio_filename)

            engine.save_to_file(extracted_text, audio_filepath)
            engine.runAndWait()
            engine.stop()  # Ensure TTS engine resets after run

            files = [f for f in os.listdir(MEDIA_FOLDER) if f.endswith(".mp3")]
            return JsonResponse({"audio_url": f"/media/{audio_filename}", "files": files})

        except Exception as e:
            return JsonResponse({"error": f"Error processing PPTX: {str(e)}"}, status=500)

    return JsonResponse({"error": "Invalid request."}, status=400)


@csrf_exempt
def delete_audio(request):
    """Delete selected audio file and update the file list."""
    if request.method == "POST":
        file_name = request.POST.get("file_name", "").strip()

        if not file_name:
            return JsonResponse({"error": "File name is required."}, status=400)

        file_path = os.path.join(MEDIA_FOLDER, file_name)

        if os.path.exists(file_path):
            try:
                os.remove(file_path)
                files = [f for f in os.listdir(MEDIA_FOLDER) if f.endswith(".mp3")]
                return JsonResponse({"message": "File deleted successfully.", "files": files})
            except Exception as e:
                return JsonResponse({"error": f"Error deleting file: {str(e)}"}, status=500)
        else:
            return JsonResponse({"error": "File not found."}, status=404)

    return JsonResponse({"error": "Invalid request."}, status=400)
